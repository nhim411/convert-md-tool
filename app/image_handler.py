"""
Image Extractor and Describer Module
Extracts images from documents and optionally describes them using AI Vision.
"""

import os
import io
import re
import base64
import logging
from pathlib import Path
from dataclasses import dataclass
from typing import List, Optional, Tuple
from PIL import Image

logger = logging.getLogger(__name__)


@dataclass
class ExtractedImage:
    """Represents an extracted image."""
    index: int
    image_data: bytes
    format: str  # 'png', 'jpeg', etc.
    source_page: Optional[int] = None
    description: Optional[str] = None
    ocr_text: Optional[str] = None


class ImageExtractor:
    """
    Extracts images from various document formats.
    """

    @staticmethod
    def extract_from_docx(file_path: str) -> List[ExtractedImage]:
        """Extract images from Word document."""
        try:
            from docx import Document
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
        except ImportError:
            logger.warning("python-docx not installed, skipping DOCX image extraction")
            return []

        images = []
        try:
            doc = Document(file_path)

            # Extract from relationships
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        ext = Path(rel.target_ref).suffix.lower().replace('.', '')
                        if ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp'):
                            images.append(ExtractedImage(
                                index=len(images) + 1,
                                image_data=image_data,
                                format=ext if ext != 'jpg' else 'jpeg'
                            ))
                    except Exception as e:
                        logger.debug(f"Failed to extract image: {e}")
        except Exception as e:
            logger.error(f"Failed to extract images from DOCX: {e}")

        return images

    @staticmethod
    def extract_from_pptx(file_path: str) -> List[ExtractedImage]:
        """Extract images from PowerPoint presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            logger.warning("python-pptx not installed, skipping PPTX image extraction")
            return []

        images = []
        try:
            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image = shape.image
                            ext = image.ext.lower()
                            if ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp'):
                                images.append(ExtractedImage(
                                    index=len(images) + 1,
                                    image_data=image.blob,
                                    format=ext if ext != 'jpg' else 'jpeg',
                                    source_page=slide_num
                                ))
                        except Exception as e:
                            logger.debug(f"Failed to extract image from slide {slide_num}: {e}")
        except Exception as e:
            logger.error(f"Failed to extract images from PPTX: {e}")

        return images

    @staticmethod
    def extract_from_pdf(file_path: str) -> List[ExtractedImage]:
        """Extract images from PDF document."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF not installed, skipping PDF image extraction")
            return []

        images = []
        try:
            doc = fitz.open(file_path)

            for page_num in range(len(doc)):
                page = doc[page_num]
                image_list = page.get_images()

                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_data = base_image["image"]
                        ext = base_image["ext"]

                        if ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp'):
                            images.append(ExtractedImage(
                                index=len(images) + 1,
                                image_data=image_data,
                                format=ext if ext != 'jpg' else 'jpeg',
                                source_page=page_num + 1
                            ))
                    except Exception as e:
                        logger.debug(f"Failed to extract image from page {page_num + 1}: {e}")

            doc.close()
        except Exception as e:
            logger.error(f"Failed to extract images from PDF: {e}")

        return images

    @classmethod
    def extract_images(cls, file_path: str) -> List[ExtractedImage]:
        """
        Extract images from a document based on its type.

        Args:
            file_path: Path to the document

        Returns:
            List of ExtractedImage objects
        """
        ext = Path(file_path).suffix.lower()

        if ext in ('.docx', '.doc'):
            return cls.extract_from_docx(file_path)
        elif ext in ('.pptx', '.ppt'):
            return cls.extract_from_pptx(file_path)
        elif ext == '.pdf':
            return cls.extract_from_pdf(file_path)
        else:
            logger.info(f"Image extraction not supported for {ext}")
            return []

    @staticmethod
    def save_images(
        images: List[ExtractedImage],
        output_dir: str,
        base_name: str
    ) -> List[str]:
        """
        Save extracted images to disk.
        Images are saved to: {output_dir}/ with simple numbered names.

        Args:
            images: List of ExtractedImage objects
            output_dir: Directory to save images (should be {base_name}_images)
            base_name: Base name (used for reference only)

        Returns:
            List of saved image paths
        """
        os.makedirs(output_dir, exist_ok=True)
        saved_paths = []

        for img in images:
            # Simple filename: image_001.png, image_002.png, etc.
            filename = f"image_{img.index:03d}.{img.format}"
            filepath = os.path.join(output_dir, filename)

            try:
                with open(filepath, 'wb') as f:
                    f.write(img.image_data)
                saved_paths.append(filepath)
                logger.debug(f"Saved image: {filepath}")
            except Exception as e:
                logger.error(f"Failed to save image {filename}: {e}")

        return saved_paths


class AIImageDescriber:
    """
    Describes images using AI Vision APIs.
    Supports OpenAI GPT-4o and Google Gemini.
    """

    def __init__(
        self,
        provider: str = "openai",
        api_key: Optional[str] = None,
        model: Optional[str] = None
    ):
        """
        Initialize AI describer.

        Args:
            provider: 'openai' or 'gemini'
            api_key: API key for the provider
            model: Model name (optional, uses default)
        """
        self.provider = provider.lower()
        self.api_key = api_key

        if self.provider == "openai":
            self.model = model or "gpt-4o-mini"
        elif self.provider == "gemini":
            self.model = model or "gemini-1.5-flash"
        else:
            raise ValueError(f"Unknown provider: {provider}")

    def describe_image(
        self,
        image_data: bytes,
        image_format: str,
        prompt: Optional[str] = None
    ) -> Optional[str]:
        """
        Get AI description of an image.

        Args:
            image_data: Raw image bytes
            image_format: Image format (png, jpeg, etc.)
            prompt: Custom prompt (optional)

        Returns:
            Description text or None if failed
        """
        if not self.api_key:
            logger.warning("No API key provided for AI image description")
            return None

        default_prompt = """Mô tả chi tiết nội dung của hình ảnh này trong 2-3 câu ngắn gọn bằng tiếng Việt.
Nếu là biểu đồ, hãy mô tả các số liệu chính.
Nếu là diagram, hãy mô tả cấu trúc và các thành phần.
Nếu có text trong ảnh, hãy trích xuất text đó."""

        prompt = prompt or default_prompt

        try:
            if self.provider == "openai":
                return self._describe_with_openai(image_data, image_format, prompt)
            elif self.provider == "gemini":
                return self._describe_with_gemini(image_data, image_format, prompt)
        except Exception as e:
            logger.error(f"AI description failed: {e}")
            return None

    def _describe_with_openai(
        self,
        image_data: bytes,
        image_format: str,
        prompt: str
    ) -> Optional[str]:
        """Use OpenAI GPT-4o Vision."""
        try:
            from openai import OpenAI
        except ImportError:
            logger.error("openai package not installed")
            return None

        client = OpenAI(api_key=self.api_key)

        # Encode image to base64
        b64_image = base64.b64encode(image_data).decode('utf-8')
        mime_type = f"image/{image_format}"

        response = client.chat.completions.create(
            model=self.model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{b64_image}",
                                "detail": "low"
                            }
                        }
                    ]
                }
            ],
            max_tokens=300
        )

        return response.choices[0].message.content

    def _describe_with_gemini(
        self,
        image_data: bytes,
        image_format: str,
        prompt: str
    ) -> Optional[str]:
        """Use Google Gemini Vision."""
        try:
            import google.generativeai as genai
        except ImportError:
            logger.error("google-generativeai package not installed")
            return None

        genai.configure(api_key=self.api_key)
        model = genai.GenerativeModel(self.model)

        # Create PIL image from bytes
        image = Image.open(io.BytesIO(image_data))

        response = model.generate_content([prompt, image])

        return response.text

    def describe_images(
        self,
        images: List[ExtractedImage],
        prompt: Optional[str] = None
    ) -> List[ExtractedImage]:
        """
        Add descriptions to a list of images.

        Args:
            images: List of ExtractedImage objects
            prompt: Custom prompt (optional)

        Returns:
            Same list with descriptions added
        """
        for img in images:
            description = self.describe_image(img.image_data, img.format, prompt)
            if description:
                img.description = description

        return images


def format_images_for_markdown(
    images: List[ExtractedImage],
    images_dir: str,
    base_name: str,
    relative_path: bool = True
) -> str:
    """
    Format extracted images as markdown text.
    Images are stored in a per-file folder: {base_name}_images/

    Args:
        images: List of ExtractedImage with optional descriptions
        images_dir: Directory where images are saved
        base_name: Base name for image files
        relative_path: Use relative paths in markdown

    Returns:
        Markdown formatted string
    """
    if not images:
        return ""

    lines = ["\n\n---\n\n## Hình ảnh trong tài liệu\n"]

    for img in images:
        # Simple filename without base_name prefix since folder already includes it
        filename = f"image_{img.index:03d}.{img.format}"

        if relative_path:
            # Per-file folder: document_images/image_001.png
            img_path = f"./{base_name}_images/{filename}"
        else:
            img_path = os.path.join(images_dir, filename)

        # Add image reference
        page_info = f" (Trang {img.source_page})" if img.source_page else ""
        lines.append(f"\n### Hình {img.index}{page_info}\n")
        lines.append(f"![Hình {img.index}]({img_path})\n")

        # Add description if available
        if img.description:
            lines.append(f"\n> **Mô tả:** {img.description}\n")

        # Add OCR text if available
        if img.ocr_text:
            lines.append(f"\n> **Text trong ảnh:**\n> {img.ocr_text}\n")

    return "\n".join(lines)
